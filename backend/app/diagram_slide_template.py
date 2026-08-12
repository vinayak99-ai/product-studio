import io

from pptx import Presentation
from pptx.util import Inches

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _save_bytes(prs: Presentation) -> bytes:
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def build_diagram_slide_pptx(png_bytes: bytes) -> bytes:
    prs = _new_presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(
        io.BytesIO(png_bytes), 0, 0, width=Inches(SLIDE_W_IN), height=Inches(SLIDE_H_IN)
    )
    return _save_bytes(prs)
