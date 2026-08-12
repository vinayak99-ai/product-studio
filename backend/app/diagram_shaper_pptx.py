"""Deterministic renderer: DiagramShaperResult -> native PPTX shapes. This is
the actual point of Diagram Shaper -- real add_shape/add_connector/textbox
calls, so every box is a movable PPT shape and every label is genuinely
editable text, not a screenshot or a glyph outline. Walks the exact same
JSON diagram_shaper_svg.py renders to SVG; the two never call into each
other, they're independent consumers of one shared source of truth.

No Playwright, no screenshot -- this tool needs no browser rendering at all.
"""

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.diagram_shaper_models import CANVAS_WIDTH, DiagramShaperResult, Shape, ShapeKind
from app.diagram_shaper_svg import ACCENT, INK, MUTED, PAPER, connector_endpoints, shape_by_id

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

_SCALE = SLIDE_W_IN / CANVAS_WIDTH  # canvas is 16:9, same ratio as the slide

_SHAPE_KIND_TO_MSO: dict[ShapeKind, int] = {
    "oval": MSO_SHAPE.OVAL,
    "rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "dot": MSO_SHAPE.OVAL,
    "cylinder": MSO_SHAPE.CAN,
    "hexagon": MSO_SHAPE.HEXAGON,
}


def _in(value: float) -> Emu:
    return Inches(value * _SCALE)


def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _save_bytes(prs: Presentation) -> bytes:
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _set_fill(shape, hex_color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))
    shape.line.fill.background()
    shape.shadow.inherit = False


def _add_shape_text(shape, text: str, *, size: int = 14) -> None:
    if not text:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string("FFFFFF")


def _add_textbox(slide, x_in: float, y_in: float, w_in: float, h_in: float, text: str, *, size: int = 11) -> None:
    box = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(INK.lstrip("#"))


def _add_shape(slide, shape: Shape) -> None:
    mso_kind = _SHAPE_KIND_TO_MSO[shape.kind]
    color = ACCENT if shape.accent else INK

    if shape.kind == "dot":
        cx, cy = shape.x + shape.w / 2, shape.y + shape.h / 2
        radius = 8
        pptx_shape = slide.shapes.add_shape(mso_kind, _in(cx - radius), _in(cy - radius), _in(2 * radius), _in(2 * radius))
        _set_fill(pptx_shape, color)
        return

    pptx_shape = slide.shapes.add_shape(mso_kind, _in(shape.x), _in(shape.y), _in(shape.w), _in(shape.h))
    _set_fill(pptx_shape, color)
    if shape.kind == "rectangle" and hasattr(pptx_shape, "adjustments"):
        pptx_shape.adjustments[0] = 0.08
    _add_shape_text(pptx_shape, shape.label)


def _add_group_box(slide, group) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _in(group.x), _in(group.y), _in(group.w), _in(group.h))
    box.fill.background()
    box.line.color.rgb = RGBColor.from_string(MUTED.lstrip("#"))
    box.line.width = Pt(1.25)
    box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    box.shadow.inherit = False
    _add_textbox(slide, group.x * _SCALE + 0.08, group.y * _SCALE + 0.05, group.w * _SCALE - 0.16, 0.25, group.label.upper(), size=9)


def _add_connector(slide, connector, shapes_by_id: dict[str, Shape]) -> None:
    source = shapes_by_id.get(connector.from_id)
    target = shapes_by_id.get(connector.to_id)
    if source is None or target is None:
        return

    (sx, sy), (ex, ey) = connector_endpoints(source, target)
    color = ACCENT if connector.accent else MUTED
    is_bent = abs(sx - ex) > 24 and abs(sy - ey) > 24
    kind = MSO_CONNECTOR.ELBOW if is_bent else MSO_CONNECTOR.STRAIGHT

    line = slide.shapes.add_connector(kind, _in(sx), _in(sy), _in(ex), _in(ey))
    line.line.color.rgb = RGBColor.from_string(color.lstrip("#"))
    line.line.width = Pt(1.5)
    line.shadow.inherit = False

    if connector.label:
        mid_x, mid_y = (sx + ex) / 2 * _SCALE, (sy + ey) / 2 * _SCALE
        _add_textbox(slide, mid_x - 0.6, mid_y - 0.18, 1.2, 0.24, connector.label, size=10)


def build_diagram_shaper_pptx(result: DiagramShaperResult) -> bytes:
    prs = _new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    _set_fill(background, PAPER)

    shapes_by_id = shape_by_id(result)

    for group in result.groups:
        _add_group_box(slide, group)
    for connector in result.connectors:
        _add_connector(slide, connector, shapes_by_id)
    for shape in result.shapes:
        _add_shape(slide, shape)

    return _save_bytes(prs)
